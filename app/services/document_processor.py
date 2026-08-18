import io
import logging
from typing import Any

import chardet
import openpyxl
import pdfplumber
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


def _for_log(value: object) -> str:
    """Make a user-supplied value safe to put in a log line.

    Upload filenames and MIME types are attacker-controlled. Without this a
    filename containing a newline can inject fabricated entries into the log,
    and an escape sequence can drive the terminal of whoever reads it.
    Ordinary values pass through unchanged.
    """
    return "".join(ch if ch.isprintable() else " " for ch in str(value))[:200]


class DocumentProcessor:
    """Service for extracting text from various document types."""

    @staticmethod
    def extract_text(file_bytes: bytes, mime_type: str, file_name: str = "") -> str | None:
        """
        Extract text from document bytes based on MIME type.
        Supports: PDF, Word, Excel, PowerPoint, and text files.
        """
        try:
            # PDF files
            if mime_type == "application/pdf":
                return DocumentProcessor._extract_from_pdf(file_bytes)

            # Word documents (.docx)
            if mime_type in [
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword",
            ]:
                return DocumentProcessor._extract_from_word(file_bytes)

            # Excel files (.xlsx, .xls)
            if mime_type in [
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
            ]:
                return DocumentProcessor._extract_from_excel(file_bytes)

            # PowerPoint files (.pptx, .ppt)
            if mime_type in [
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint",
            ]:
                return DocumentProcessor._extract_from_powerpoint(file_bytes)

            # Text files
            if mime_type.startswith("text/"):
                return DocumentProcessor._extract_from_text(file_bytes)

            # Try to detect and process as text if unknown
            logger.warning(f"Unknown MIME type {_for_log(mime_type)}, attempting text extraction")
            return DocumentProcessor._extract_from_text(file_bytes)

        except Exception as e:
            logger.error(
                f"Error extracting text from {_for_log(file_name)} (type: {_for_log(mime_type)}): {e}"
            )
            return None

    @staticmethod
    def _extract_from_pdf(pdf_bytes: bytes) -> str | None:
        """Extract text from PDF files."""
        try:
            # Try pdfplumber first (better text extraction)
            with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
                text_parts = []
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                return "\n\n".join(text_parts)
        except Exception as e:
            logger.warning(f"pdfplumber extraction failed: {e}, trying PyPDF2")
            try:
                # Fallback to PyPDF2
                pdf_file = io.BytesIO(pdf_bytes)
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                text_parts = []
                for page in pdf_reader.pages:  # type: ignore[assignment]  # PyPDF2 fallback page type
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                return "\n\n".join(text_parts)
            except Exception as e2:
                logger.error(f"PyPDF2 extraction also failed: {e2}")
                return None

    @staticmethod
    def _extract_from_word(docx_bytes: bytes) -> str | None:
        """Extract text from Word documents (.docx)."""
        try:
            doc = DocxDocument(io.BytesIO(docx_bytes))
            text_parts = []

            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        text_parts.append(row_text)

            return "\n\n".join(text_parts) if text_parts else None
        except Exception as e:
            logger.error(f"Error extracting from Word document: {e}")
            return None

    @staticmethod
    def _extract_from_excel(xlsx_bytes: bytes) -> str | None:
        """Extract text from Excel files (.xlsx)."""
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), data_only=True)
            text_parts = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"\n--- Sheet: {sheet_name} ---\n")

                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text_parts.append(row_text)

            return "\n".join(text_parts) if text_parts else None
        except Exception as e:
            logger.error(f"Error extracting from Excel file: {e}")
            return None

    @staticmethod
    def _extract_from_powerpoint(pptx_bytes: bytes) -> str | None:
        """Extract text from PowerPoint files (.pptx)."""
        try:
            prs = Presentation(io.BytesIO(pptx_bytes))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n--- Slide {slide_num} ---\n")

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

            return "\n".join(text_parts) if text_parts else None
        except Exception as e:
            logger.error(f"Error extracting from PowerPoint file: {e}")
            return None

    @staticmethod
    def _extract_from_text(text_bytes: bytes) -> str | None:
        """Extract text from text files (with encoding detection)."""
        try:
            # Detect encoding
            detected = chardet.detect(text_bytes)
            encoding = detected.get("encoding", "utf-8")

            # Try to decode
            try:
                return text_bytes.decode(encoding)  # type: ignore[arg-type]  # None is handled below
            except (UnicodeDecodeError, LookupError):
                # Fallback to utf-8 with error handling
                return text_bytes.decode("utf-8", errors="ignore")
        except Exception as e:
            logger.error(f"Error extracting from text file: {e}")
            return None

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[dict[str, Any]]:
        """
        Split text into chunks with overlap.
        Returns list of dicts with 'text', 'index', 'start_char', and 'end_char' keys.
        """
        if not text:
            return []

        chunks = []
        start = 0
        index = 0

        while start < len(text):
            end = start + chunk_size
            chunk_text = text[start:end]

            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings
                last_period = chunk_text.rfind(".")
                last_newline = chunk_text.rfind("\n")
                break_point = max(last_period, last_newline)

                if break_point > chunk_size * 0.5:  # Only break if we're at least halfway
                    chunk_text = chunk_text[: break_point + 1]
                    end = start + break_point + 1

            chunks.append({"text": chunk_text.strip(), "index": index, "start_char": start, "end_char": end})

            start = end - chunk_overlap  # Overlap for context
            index += 1

        return chunks

    @staticmethod
    def detect_sensitive_data(text: str) -> list[str]:
        """
        Detect sensitive data (PII, secrets) in text.
        Returns a list of warning messages with the detected content context.
        """
        import re

        warnings = []

        # Patterns for sensitive data
        patterns = {
            "Email Address": r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}",
            "Credit Card Number": r"\b(?:\d{4}[-\s]?){3}\d{4}\b",
            "Social Security Number (SSN)": r"\b\d{3}-\d{2}-\d{4}\b",
            "API Key / Secret": r"(?i)(?:api[_-]?key|secret[_-]?key|access[_-]?token)[\s:=]+([a-zA-Z0-9_\-]{20,})",  # noqa: E501
            # "Phone Number": r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b' # Too many false positives often
        }

        lines = text.split("\n")
        for i, line in enumerate(lines):
            for label, pattern in patterns.items():
                matches = re.finditer(pattern, line)
                for match in matches:
                    # Provide snippet context
                    start = max(0, match.start() - 20)
                    end = min(len(line), match.end() + 20)
                    snippet = line[start:end].strip()
                    warnings.append(f"Line {i + 1}: Potential {label} detected: '...{snippet}...'")

        # Limit total warnings to avoid overwhelming response
        return warnings[:10]
